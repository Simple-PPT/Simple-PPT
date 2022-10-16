import os
import platform
import sys
import webbrowser
from subprocess import Popen

import cmake
import docx
import fake_useragent
from bs4 import BeautifulSoup
from PySide6.QtCore import *
from PySide6.QtWidgets import *
from pptx import Presentation
import pypandoc
from requests import *
from ui_NoDone import Ui_MainWindow


class MainWindow(QMainWindow, Ui_MainWindow):
    def __init__(self, parent=None):
        super(MainWindow, self).__init__(parent)
        self.filepath = ""
        self.filetype = "docx"
        self.setupUi(self)
        self.trans = QTranslator()
        self.version = 0.1
        self.session = session()
        self.ua = fake_useragent.UserAgent()
        self.arch, self.exe = platform.architecture()

    def ShowAbout(self):
        QMessageBox.about(self, QCoreApplication.translate("about", "About Simple PPT 0.1", None), QCoreApplication.translate("about", "<h2>Simple PPT</h2>version:0.0.2<br/>It can make a PPT for you.<br/>More see:<a href=https://github.com/Simple-PPT/Simple-PPT>Github</a>", None))

    def OpenFeedback(self):
        webbrowser.open_new("https://github.com/Simple-PPT/Simple-PPT/issues/new")

    def checkupdate(self):
        newversion = self.session.get('https://api.github.com/repos/Simple-PPT/Simple-PPT/releases/latest').json()
        tagnewversion = newversion["tag_name"]
        name = newversion["name"]
        try:
            tagname = float(newversion["tag_name"])
        except Exception as e:
            QMessageBox.critical(self, QCoreApplication.translate("MessageBox", "Error", None), QCoreApplication.translate("MessageBox", "Error:<br/>%s<br/><br/><a href=https://github.com/Simple-PPT/Simple-PPT/issues/new>Feedback Error</a>", None)%e, QMessageBox.Ok)
        else:
            if tagname > self.version:
                to_update = QMessageBox.question(self, QCoreApplication.translate("MessageBox", "New version", None), QCoreApplication.translate("MessageBox", "Find a new version:%s.\nWould you like to update?", None)% name, buttons=QMessageBox.Yes|QMessageBox.No, defaultButton=QMessageBox.Yes)
                if to_update == QMessageBox.Yes:
                    print(self.arch)
                    if self.arch == '64bit':
                        try:
                            req = self.session.get("https://github.com/Simple-PPT/Simple-PPT/releases/download/%s/Simple-PPT_%s_x64_Setup.exe"%(tagnewversion, name), headers={"User-Agent":self.ua.random}, stream=True)
                        except exceptions.ConnectTimeout or exceptions.ReadTimeout or exceptions.ConnectionError:
                            QMessageBox.warning(self, QCoreApplication.translate("MessageBox", "Download time out", None), QCoreApplication.translate("MessageBox", "Download time out.</br>Please try again.", None), QMessageBox.Ok)
                        else:
                            if req.status_code == 200:
                                downloadpath = os.path.basename("https://github.com/Simple-PPT/Simple-PPT/releases/download/%s/Simple-PPT_%s_x64_Setup.exe"%(tagnewversion, name))
                                with open(downloadpath, 'wb') as f:
                                    for i in req.iter_content(1024):
                                        if i:
                                            f.write(i)
                                downloadfilefilepath = os.path.realpath(downloadpath)
                                DFFP = downloadfilefilepath
                                Popen(DFFP)
                            else:
                                QMessageBox.critical(self, QCoreApplication.translate("MessageBox", "Error", None), QCoreApplication.translate("MessageBox", "Error:<br/>We don't know why, error code %s.<br/><br/><a href=https://github.com/Simple-PPT/Simple-PPT/issues/new>Feedback Error</a>", None)%req.status_code, QMessageBox.Ok)

                    elif self.arch == '32bit':
                        try:
                            req = self.session.get("https://github.com/Simple-PPT/Simple-PPT/releases/download/%s/Simple-PPT_%s_x32_Setup.exe"%(tagnewversion, name), headers={'User-Agent':self.ua.random}, stream=True)
                        except exceptions.ConnectTimeout or exceptions.ReadTimeout or exceptions.ConnectionError:
                            QMessageBox.warning(self, QCoreApplication.translate("MessageBox", "Download time out", None), QCoreApplication.translate("MessageBox", "Download time out.</br>Please try again.", None), QMessageBox.Ok)
                        else:
                            if req.status_code == 200:
                                downloadpath = os.path.basename("https://github.com/Simple-PPT/Simple-PPT/releases/download/%s/Simple-PPT_%s_x32_Setup.exe"%(tagnewversion, name))
                                with open(downloadpath, 'wb') as f:
                                    for i in req.iter_content(1024):
                                        if i:
                                            f.write(i)
                                downloadfilefilepath = os.path.realpath(downloadpath)
                                DFFP = downloadfilefilepath
                                Popen(DFFP)
                            else:
                                QMessageBox.critical(self, QCoreApplication.translate("MessageBox", "Error", None), QCoreApplication.translate("MessageBox", "Error:<br/>We don't know why, error code %s.<br/><br/><a href=https://github.com/Simple-PPT/Simple-PPT/issues/new>Feedback Error</a>", None)%req.status_code, QMessageBox.Ok)
            else:
                QMessageBox.information(self, QCoreApplication.translate("MessageBox", "No new version", None), QCoreApplication.translate("MessageBox", "Your version is latest."), QMessageBox.Yes)
    def chagefile(self):
        filePath, filetype = QFileDialog.getOpenFileName(self, QCoreApplication.translate("QFileDialog", "Open a file", None), r"c:\\", QCoreApplication.translate("QFileDialog", "Documentation(*.docx *.md)", None))
        if filePath == None or filePath == "":
            pass
        else:
            if os.path.isfile(filePath):
                if filePath.endswith(".docx"):
                    self.search_filePath([filePath, "docx"])
                if filePath.endswith(".md"):
                    pypandoc.convert_file(filePath, 'docx', 'md', outputfile='./file1.docx')
            else:
                QMessageBox.warning(self, QCoreApplication.translate("MessageBox", "wrong file path", None), QCoreApplication.translate("MessageBox","The file's path is wrong!", None), QMessageBox.Yes)

    def CreatPPT(self):
        ppt = Presentation()
        titlestyle = self.comboBox_2.currentIndex()
        textstyle = self.comboBox_4.currentIndex()
        title_slide = ppt.slides.add_slide(ppt.slide_layouts[titlestyle])
        title = title_slide.shapes.title
        title.text = self.lineEdit.text()
        subtitle = title_slide.placeholders[1]
        subtitle.text = self.lineEdit_2.text()

    def search_filePath(self, *args):
        self.filepath = args[0]
        self.filetype = args[1]

    def translate_to_en_UK(self):
        self.trans.load("./i18n/Eng_UK")
        _app = QApplication.instance()
        _app.installTranslator(self.trans)
        self.retranslateUi(self)

    def translate_to_zh_hans(self):
        self.trans.load("./i18n/zh_Hans_CN")
        _app = QApplication.instance()
        _app.installTranslator(self.trans)
        self.retranslateUi(self)

    def translate_to_zh_hant(self):
        self.trans.load("./i18n/zh_Hant_CN")
        _app = QApplication.instance()
        _app.installTranslator(self.trans)
        self.retranslateUi(self)


def main():
    app = QApplication()
    win = MainWindow()
    win.show()
    sys.exit(app.exec())

main()
