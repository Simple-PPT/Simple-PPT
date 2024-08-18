from pptx import Presentation
from docx import Document
from pathlib import Path
import pypandoc

class PPT:
    def __init__(self, ppt_name: str):
        self.ppt_name = ppt_name
        self.now_page = 0
        self.all_pages = 0
        self.slide = 0
        self.text = ""
        self.power_point = Presentation()
        self.save()
    
    def del_page(self):
        sli = list(self.power_point.slides._sldIdLst)
        self.power_point.slides._sldIdLst.remove(sli[-self.now_page])
        self.save()
    
    def page(self, idx: int, slide = 0):
        self.now_page = idx
        if self.now_page > self.all_pages:
            self.all_pages = self.now_page
        self.allpages(self.now_page, slide)
        self.save()
    
    def get_ppt_pages(self):
        total_page_count = len(self.power_point.slides)
        return total_page_count
    
    def add_slide(self):
        self.power_point.slides.add_slide(self.power_point.slide_layouts[self.slide])
    
    def write_str(self, string: str):
        self.power_point.slides[self.now_page - 1].shapes.placeholders[1].text = string
        self.save()
    
    def append_str(self, string: str):
        self.power_point.slides[self.now_page - 1].shapes.placeholders[1].text_frame.add_paragraph().text = string
        self.save()
    
    def write_docx(self, docp: str | Path):
        doc = Document(docp)
        place = self.power_point.slides[self.now_page - 1].shapes.placeholders[1]
        place.text = ""
        for paragraph in doc.paragraphs:
            place.text_frame.add_paragraph().text = paragraph.text
        self.save()
    
    def append_docx(self, docp: str | Path):
        doc = Document(str(docp))
        place = self.power_point.slides[self.now_page - 1].shapes.placeholders[1]
        for paragraph in doc.paragraphs:
            place.text_frame.add_paragraph().text = paragraph.text
        self.save()
    
    def write_markdown(self, mkdown: str):
        tmp = Path("tmp.docx")
        pypandoc.convert_file(Path(mkdown).resolve(), "docx", "md", outputfile=tmp)
        self.write_docx(tmp)
        tmp.unlink()
        self.save()
    
    def append_markdown(self, mkdown: str):
        tmp = Path("tmp.docx")
        pypandoc.convert_file(Path(mkdown).resolve(), "docx", "md", outputfile=tmp)
        self.append_docx(tmp)
        tmp.unlink()
        self.save()
    
    def allpages(self, idx: int, slide = 0):
        self.all_pages = idx
        while self.all_pages > self.get_ppt_pages():
            self.slide = slide
            self.add_slide()
        self.save()
    
    def save(self):
        self.power_point.save(self.ppt_name+".pptx")
