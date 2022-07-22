print("Note:Before make a PPT(.pptx) file,please make the info-user.xlsx file first.")

from tkinter import *
from tkinter import ttk
from pptx import Presentation
import pandas as pd
import time,warnings
from tqdm import tqdm
import numpy as np
from openpyxl import *
from sys import exit

exda = pd.DataFrame(
    {
    'pagenumbers(num)':1,
    'style(1 - 11)':1,
    'PPT\'s name(string)':'write',
    'title(string)':'write',
    'subtitle(string)':'write'
    },
    index = [0]
)
exda.to_excel('info.xlsx',sheet_name = 'read',index = False)
class wait(object):
    def to_wait(self):
        self.wait_time = input("Please input your wait time for write info.xlsx(s):")
        try:
            print(pd.read_excel('./info.xlsx','read'))
            for i in tqdm(range(int(self.wait_time))):
                time.sleep(1)
        except:
            warnings.showwarning(f"You can't input {type(self.wait_time)}!",category = UserWarning,lineno = 26,filename = 'b1.py')

def make_PPT(style,name,title_input,subtitle_input):
    ppt =Presentation()
    to_int_style = np.array(style).item()
    title_slide = ppt.slides.add_slide(ppt.slide_layouts[to_int_style - 1])
    title = title_slide.shapes.title
    title.text = title_input
    subtitle = title_slide.placeholders[1]
    subtitle.text = subtitle_input
    ppt.save(str(name) + ".pptx")
wait().to_wait()
incase_OK = input("Are you write info.xlsx yet? (y/n):")
while True:
    def scaner():
        loadwb = load_workbook('./info.xlsx')
        loadsheet = loadwb['read']
        make_PPT(style = int((loadsheet.cell(row = 2,column = 2)).value),name = str((loadsheet.cell(row = 2,column = 3)).value),title_input=(loadsheet.cell(row = 2,column = 4)).value,subtitle_input=(loadsheet.cell(row = 2,column = 5)).value)
    if incase_OK == "y" or incase_OK == "Y":
        scaner()
        exit()

    elif incase_OK == "n" or incase_OK == "N":
        wait().to_wait()
        incase_OK = input("Are you write info.xlsx yet? (y/n):")
        scaner()
        exit()

    else:
        print('Only y or n')
        wait().to_wait()
        incase_OK = input("Are you write info.xlsx yet? (y/n):")
        scaner()
        exit()
