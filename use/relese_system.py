from pptx import Presentation
import pandas as pd
import numpy as np
import openpyxl
import json,os
zh_Hans_CN_alls = ['语言','版本','请先输入写info.xlsx文件的时间（秒）：','好了','检查更新','错误','不要输入文字！','时间到','你写好info.xlsx文件了吗？','完成！','您的PPT完成了！','提示','\'github.com\'没有响应','新的','有新的版本！','下载','你想下载新的版本吗？','你的版本是最新的！']
en_US_alls = ['Language','Version','Please input write the info.xlsx file\'s time frist(s):','OK','Check new','Error','Don\'t input string!','Timed up','Have you finished writing the info.xlsx file?','Done!','Your PPT is OK!','info','\'github.com\' has no response','new','Have new version!','Download','Do you want download new version?','Your version is new!']
ah = {}
data = {"language":False}
class system(object):
    '''
    Help for GUI.
    '''
    def make_excel(self):
        '''
        Make a excel file for GUI.
        '''
        self.exda = pd.DataFrame(
            {
            'pagenumbers(num)':1,
            'style(1 - 11)':1,
            'PPT\'s name(string)':'write',
            'title(string)':'write',
            'subtitle(string)':'write',
            'text style(1 - 11)': 1,
            'Do you have text title(Yes/No)':'No',
            'text title(optional)(string)':'write',
            'text(string)':'write'
            },
            index = [0]
        )
        self.exda.to_excel('info.xlsx',sheet_name = 'read',index = False)
    def make_PPT(self,many:int,style:int,name:str,title_input:str,subtitle_input:str,every_style:int,have_title : bool or str,title_E_text:str,word:str):
        ppt = Presentation()
        to_int_style = np.array(style).item()
        title_slide = ppt.slides.add_slide(ppt.slide_layouts[to_int_style - 1])
        title = title_slide.shapes.title
        title.text = title_input
        subtitle = title_slide.placeholders[1]
        subtitle.text = subtitle_input
        for i in range(1,many):
            to_int_ES = np.array(every_style).item()
            slide = ppt.slides.add_slide(ppt.slide_layouts[to_int_ES - 1])
            Etitle = slide.shapes.title
            if have_title == "Yes":
                Etitle.text = title_E_text
            Eword = slide.placeholders[1]
            Eword.text = word
        ppt.save(str(name) + ".pptx")

    def read_and_make(self):
        self.loadwb = openpyxl.load_workbook("./info.xlsx")
        self.loadsheet = self.loadwb['read']
        self.make_PPT(many = int(self.loadsheet.cell(2,1).value),style = int((self.loadsheet.cell(row = 2,column = 2)).value),name = str((self.loadsheet.cell(row = 2,column = 3)).value),title_input=(self.loadsheet.cell(row = 2,column = 4)).value,subtitle_input=(self.loadsheet.cell(row = 2,column = 5)).value,every_style = int(self.loadsheet.cell(row = 2,column = 6).value),have_title = (self.loadsheet.cell(2,7)).value,title_E_text= str((self.loadsheet.cell(2,8)).value),word = str((self.loadsheet.cell(2,9)).value))

    def load_language(self):
        global ah,data
        if os.path.isfile('./info_server.json'):
            with open('./info_server.json') as f:
                f = json.load(f)
                ah = json.dumps(f)
                ah = json.loads(ah)
            if ah["language"] == "en-US":
                alls = en_US_alls
                return alls
            elif ah["language"] == "zh-Hans-CN":
                alls = zh_Hans_CN_alls
                return alls
            else:
                os.remove('./info_server.json')
        else:
            alls = en_US_alls
            return alls

    def write_en_US(self):
        global data
        with open('info_server.json','w') as fp:
            data["language"] = "en-US"
            json.dump(data,fp,indent = 2)

    def write_zh_Hans_CN(self):
        global data
        with open('./info_server.json','w') as fp:
            data["language"] = "zh-Hans-CN"
            json.dump(data,fp,indent = 2)
