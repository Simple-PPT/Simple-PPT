print("Note:Before make a PPT(.pptx) file,please make the info-user.xlsx file first.")
print("Note:info-user.xlsx should have a sheet,the sheet's name is'user'")

from pptx import Presentation
import pandas as pd
import time
from tqdm import tqdm
import numpy as np
from openpyxl import *

write_nobugs = Workbook()
ws = write_nobugs.create_sheet('read')
exda =[
        ['pagenumbers(num)','style(1 - 11)','PPT\'s name(string)','title(string)','subtitle(string)'],
        ['write','write','write','write','write']
    ]
for ed in exda:
    ws.append(ed)
del write_nobugs['Sheet']
write_nobugs.save('info.xlsx')
class wait(object):
    def to_wait(self):
        self.wait_time = input("Please input your wait time for write info.xlsx(s):")
        try:
            print(pd.read_excel('./info.xlsx','read'))
            for i in tqdm(range(int(self.wait_time))):
                time.sleep(1)
        except Exception as e:
            print(e)

def make_PPT(style,name,title_input,subtitle_input):
    ppt =Presentation()
    to_int_style = np.array(style).item()
    title_slide = ppt.slides.add_slide(ppt.slide_layouts[to_int_style - 1])
    title = title_slide.shapes.title
    title.text = title_input
    subtitle = title_slide.placeholders[1]
    subtitle.text = subtitle_input
    ppt.save(str(name) + ".pptx")
wait().to_wait()
incase_OK = input("Are you write info.xlsx yet? (y/n):")
while True:
    def scaner():
        if incase_OK == "y" or incase_OK  == "Y":
            loadwb = load_workbook('./info-user.xlsx')
            loadsheet = loadwb['Sheet1']
            make_PPT(style = int((loadsheet.cell(row = 2,column = 2)).value),name = str((loadsheet.cell(row = 2,column = 3)).value),title_input=(loadsheet.cell(row = 2,column = 4)).value,subtitle_input=(loadsheet.cell(row = 2,column = 5)).value)
    if incase_OK == "n" or incase_OK == "N":
        wait().to_wait()
        incase_OK = input("Are you write info.xlsx yet? (y/n):")
        scaner()
    else:
        print('Only y or n')
        wait().to_wait()
        incase_OK = input("Are you write info.xlsx yet? (y/n):")
        scaner()
